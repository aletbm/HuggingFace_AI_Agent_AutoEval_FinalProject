
import os
import requests
import tempfile
from langchain.docstore.document import Document
from langchain.tools import Tool, StructuredTool, tool
from typing import Tuple, List
from langchain_core.messages import AIMessage
from langchain_community.tools import DuckDuckGoSearchResults, DuckDuckGoSearchRun
from langchain_experimental.tools.python.tool import PythonREPLTool
from langchain_community.agent_toolkits.playwright.toolkit import PlayWrightBrowserToolkit
from langchain_community.tools.playwright.utils import create_async_playwright_browser
from langchain_community.document_loaders.wikipedia import WikipediaLoader
from langchain_community.document_loaders import WebBaseLoader, PyPDFLoader
from langchain_community.document_loaders.arxiv import ArxivLoader
from langchain_community.document_loaders.text import TextLoader
from langchain_community.document_loaders.word_document import Docx2txtLoader
import yt_dlp
import json
from Bio.PDB import PDBParser
import pandas as pd
from openpyxl import load_workbook
import math
import pdfplumber
from pptx import Presentation
import zipfile
from ratelimit import limits, sleep_and_retry
from tenacity import retry, wait_exponential, stop_after_attempt, retry_if_exception_type
from duckduckgo_search.exceptions import DuckDuckGoSearchException
from transformers import Qwen2_5_VLForConditionalGeneration, AutoProcessor
from qwen_vl_utils import process_vision_info
import whisper

model_whisper = whisper.load_model("base")
model_image = Qwen2_5_VLForConditionalGeneration.from_pretrained(
    "Qwen/Qwen2.5-VL-3B-Instruct",
    torch_dtype="auto",
    device_map="auto"
)
processor_image = AutoProcessor.from_pretrained("Qwen/Qwen2.5-VL-3B-Instruct")

@tool
def add(a: str, b: str) -> str:
    """Adds two numbers provided as strings and returns the result as a string."""
    return str(float(a) + float(b))

@tool
def subtract(a: str, b: str) -> str:
    """Subtracts the second number from the first, both provided as strings, and returns the result as a string."""
    return str(float(a) - float(b))

@tool
def multiply(a: str, b: str) -> str:
    """Multiplies two numbers provided as strings and returns the result as a string."""
    return str(float(a) * float(b))

@tool
def divide(a: str, b: str) -> str:
    """Divides the first number by the second, both provided as strings, and returns the result as a string. The divisor must not be zero."""
    return str(float(a) / float(b))

@tool
def power(a: str, b: str) -> str:
    """Raises the first number (base) to the power of the second number (exponent), both provided as strings, and returns the result as a string."""
    return str(float(a) ** float(b))

@tool
def square_root(a: str) -> str:
    """Calculates the square root of a number provided as a string and returns the result as a string. The number must be non-negative."""
    return str(math.sqrt(float(a)))

@tool
def get_information_from_wikipedia(query: str) -> str:
    """
    Search for relevant Wikipedia pages based on a user query and return their URLs.

    This tool uses WikipediaLoader to retrieve up to 10 Wikipedia articles related to the input query.
    The output consists of a formatted list of URLs pointing to these articles. These URLs can be used
    by web automation tools or agents to navigate, scrape, and extract valuable information such as
    text, tables, infoboxes, images, and references that may help answer the user's question.

    Args:
        query (str): A user-provided search query.

    Returns:
        str: A formatted explanation with a list of Wikipedia URLs related to the query.
    """
    search_docs = WikipediaLoader(query=query, load_max_docs=10).load()
    urls = [doc.metadata["source"] for doc in search_docs]
    url_list = "\n".join(f"- {url}" for url in urls)

    return f"""Wikipedia search completed for query: "{query}".

You must now browse the following Wikipedia pages to extract detailed information. Use `navigate_browser` to explore each URL:

Relevant Wikipedia URLs:
{url_list}

Instructions:
- Visit each URL using `navigate_browser`.
- Extract full page content, infoboxes, tables, and references.
- Use hyperlinks from the page to recursively follow relevant internal links if necessary."""

@tool
def get_information_from_arxiv(query: str) -> str:
    """
    Searches the arXiv database for academic papers related to a given query and returns up to 3 relevant results,
    including source metadata and a preview of each paper's content.

    This function queries the arXiv database for papers matching the provided search term. It retrieves up to 3 papers
    with content limited to 10,000 characters each. The results are formatted in XML-like structure with metadata
    such as the title, authors, and page number (if available).

    Args:
        query (str): The search term or topic to query on arXiv (e.g., "machine learning").

    Returns:
        str: A formatted string containing the metadata (title, authors, page number) and content preview
             for each of the top 3 relevant papers found on arXiv.
    """
    search_docs = ArxivLoader(query=query, load_max_docs=3, doc_content_chars_max=10000).load()
    content = ""
    for doc in search_docs:
        content += f"Title: {doc.metadata['Title']}\n"
        content += f"Authors: {doc.metadata['Authors']}\n"
        content += f"Summary: {doc.metadata['Summary']}\n"
        content += f"Published: {doc.metadata['Published']}\n\n"

    return f"""Arxiv search completed for query: "{query}".

Relevant Arxiv papers:
{content}

Instructions:
- Search the paper link throught it's title.
- Visit each URL using `navigate_browser`.
- Extract full page content, infoboxes, tables, and references.
- Use hyperlinks from the page to recursively follow relevant internal links if necessary."""

@tool
def get_web_page(url: str) -> str:
    """
    Use this to extract text from basic HTML pages. Works well for static websites without JavaScript.

    This function retrieves the content of a webpage at the specified URL and extracts all the visible text.

    Args:
        url (str): The URL of the public webpage to fetch (e.g., "https://example.com").

    Returns:
        str: A string containing the visible text content extracted from the webpage.

    Example:
        url = "https://en.wikipedia.org/wiki/Python_(programming_language)"
        page_text = get_web_page(url)
        print(page_text)

    In this example, the function will return the visible text from the Wikipedia page about "Python programming language",
    excluding JavaScript-generated or hidden content, formatted as plain text.
    """
    loader = WebBaseLoader(url)
    docs = loader.load()
    content = "\n\n".join([doc.page_content for doc in docs])
    return f"""The web page content is:

{content}

Instructions:
If the information is not sufficient:
- Visit the URL {url} using `navigate_browser`.
- Extract full page content, infoboxes, tables, and references.
- Use hyperlinks from the page to recursively follow relevant internal links if necessary."""

def download_video_from_youtube(url, tempdirname):
    """
    Downloads the best quality video from YouTube (MP4 format) and saves it in the specified directory.

    This function fetches the video from the provided YouTube URL, ensuring that the video format is MP4 with a
    resolution of at least 1080p. It also downloads subtitles if available, saving them in English as SRT files.

    Args:
        url (str): The YouTube video URL to download (e.g., "https://www.youtube.com/watch?v=example").
        tempdirname (str): The directory where the downloaded video and subtitles will be saved.

    Returns:
        str: The path to the downloaded video file (MP4 format).

    Example:
        url = "https://www.youtube.com/watch?v=example"
        tempdirname = "/tmp/videos"
        video_filepath = download_video_from_youtube(url, tempdirname)
        print(video_filepath)

    In this example, the function will download the best quality video and subtitles from the given YouTube URL
    and save them in the specified directory.
    """

    ydl_opts = {
        'format': '(bestvideo[width>=1080][ext=mp4]/bestvideo)+bestaudio/best', #Ensures best settings
        'writesubtitles': True, #Adds a subtitles file if it exists
        'writeautomaticsub': True, #Adds auto-generated subtitles file
        'subtitle': '--write-sub --sub-lang en', #writes subtitles file in english
        'subtitlesformat':'srt', #writes the subtitles file in "srt" or "ass/srt/best"
        'skip_download': False, #skips downloading the video file
        "merge_output_format": "mp4",
        'outtmpl':f"{tempdirname}/%(title)s.%(ext)s",
        'quiet': True,
        'cookiefile': "./cookies.txt"
        }
    with yt_dlp.YoutubeDL(ydl_opts) as ydl:
        info_dict = ydl.extract_info(url, download=True)
        video_title = info_dict.get('title', None)
        filepath = ydl.prepare_filename(info_dict)
    print(f"Download {filepath.split('/')[-1]} Successful!")
    return filepath

def download_audio_from_youtube(url, tempdirname):
    """
    Downloads the best quality audio (MP3 format) from a YouTube video and saves it in the specified directory.

    This function fetches the audio from the provided YouTube URL in the best available format and saves it
    as an MP3 file. It does not download the video.

    Args:
        url (str): The YouTube video URL to extract audio from (e.g., "https://www.youtube.com/watch?v=example").
        tempdirname (str): The directory where the audio file will be saved.

    Returns:
        str: The path to the downloaded audio file (MP3 format).

    Example:
        url = "https://www.youtube.com/watch?v=example"
        tempdirname = "/tmp/audio"
        audio_filepath = download_audio_from_youtube(url, tempdirname)
        print(audio_filepath)

    In this example, the function will download the best quality audio (MP3) from the YouTube URL
    and save it in the specified directory.
    """

    ydl_opts = {
        'extract_audio': True,
        'format': 'bestaudio',
        'outtmpl': f"{tempdirname}/%(title)s.mp3",
        'quiet': True,
        'cookiefile': "./cookies.txt"
        }
    with yt_dlp.YoutubeDL(ydl_opts) as ydl:
        info_dict = ydl.extract_info(url, download=True)
        audio_title = info_dict.get('title', None)
        filepath = ydl.prepare_filename(info_dict)
    print(f"Download {filepath.split('/')[-1]} Successful!")
    return filepath

def download_youtube_data(url: str) -> str:
    """
    Downloads the description and comments of a YouTube video using the yt-dlp Python API.

    Args:
        url (str): The URL of the YouTube video.

    Returns:
        str: File paths and summary of downloaded content.
    """
    ydl_opts = {
        'quiet': True,
        'extract_flat': False,
        'skip_download': True,
        'getcomments': True,
    }

    try:
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(url, download=False)

        # Guardar descripción
        description = info.get("description", "No description found.")

        # Guardar comentarios
        comments = info.get("comments", [])
        comment_texts = [c["text"] for c in comments if "text" in c]

        return (
            f"Downloaded video data:\n\n"
            f"Description: {description}\n"
            f"Comments: {comment_texts}\n"
        )

    except Exception as e:
        return f"Error extracting video data: {str(e)}"

@tool
def get_information_from_youtube(url: str) -> str:
    """
    Downloads a YouTube video, extracts the audio, and returns a transcription of the audio.

    This function first downloads the video and audio from the given YouTube URL. Then, it transcribes the
    speech from the audio file using a transcription service.

    Args:
        url (str): The YouTube video URL to download and transcribe (e.g., "https://www.youtube.com/watch?v=example").

    Returns:
        str: The transcription text of the audio content extracted from the YouTube video.

    Example:
        url = "https://www.youtube.com/watch?v=example"
        transcription = get_information_from_youtube(url)
        print(transcription)

    In this example, the function will download the video and audio from YouTube, extract the audio, and
    return a transcription of the spoken content in the video.
    """

    tempdirname = tempfile.TemporaryDirectory()
    video_filepath = download_video_from_youtube(url, tempdirname)
    audio_filepath = download_audio_from_youtube(url, tempdirname)
    youtube_data = download_youtube_data(url)

    return youtube_data + get_information_from_audio.invoke(audio_filepath)

@tool
def get_information_from_audio(filepath: str) -> str:
    """
    Transcribes speech from an audio file into text using a speech-to-text model.

    This function takes an audio file and converts the spoken content into written text using a speech-to-text model.
    It is useful when you need to extract information from voice recordings.

    Args:
        filepath (str): The path to the audio file (e.g., ".wav", ".mp3", etc.).
                         The file should contain spoken content to be transcribed.

    Returns:
        str: The transcribed text from the audio file.

    Example:
        filepath = "/path/to/audiofile.wav"
        transcription = get_information_from_audio(filepath)
        print(transcription)

    In this example, the function will transcribe the speech from the specified audio file and return the
    transcribed text that can be further analyzed or used in other processes.
    """
    result = model_whisper.transcribe(filepath)
    return (f"Transcription of audio file:\n\n{result['text']}"
            "Please verify this information by cross-checking with reliable external sources such as web searches, Wikipedia, or other knowledge bases before finalizing your response."
            "If necessary, supplement the transcription with additional relevant information to ensure completeness and accuracy.")

@tool
def get_information_from_pdf(file_path: str) -> str:
    """
    Extracts structured text content from a PDF file.

    This function reads the content of a PDF document located at the specified file path and extracts the text
    from each page. The text is returned as a single string, with the order of pages preserved.

    Args:
        file_path (str): The local path to the PDF file to be processed (e.g., "/path/to/document.pdf").

    Returns:
        str: A single string containing the combined text extracted from all pages of the PDF.
             The text is ordered by the pages as they appear in the document.

    Example:
        file_path = "/path/to/document.pdf"
        extracted_text = get_information_from_pdf(file_path)
        print(extracted_text)

    In this example, the function will extract all text content from the PDF document,
    maintaining the page order, and return it as one continuous string.
    """

    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            content = page.extract_tables()
            if content and len(content) > 0:
                print("Tables were detected in the PDF.")
            else:
                print("No tables were detected in the PDF..")
                loader = PyPDFLoader(file_path = file_path)
                documents = loader.load()
                content = '\n'.join([doc.page_content for doc in documents])

    return f"""The PDF file content:

{content}

Use this data to answer the question or perform any required analysis.
Remember, you can use all available tools to supplement with additional relevant information."""

@tool
def get_information_from_csv(file_path: str) -> str:
    """
    Loads a CSV file and returns a structured summary of its tabular content.

    This function reads the entire CSV file and returns a textual representation of the data
    intended for use by AI agents or downstream systems that operate on text-based tables.

    Args:
        file_path (str): The path to the CSV file to load.

    Returns:
        str: A formatted string containing the CSV content as a table under the label `data_table`,
             along with metadata such as the number of rows and columns.
    """
    df = pd.read_csv(file_path)

    return f"""You are given a table extracted from a CSV file with {df.shape[0]} rows and {df.shape[1]} columns.

You are given a table extracted from a CSV file.

- Columns: {df.columns.tolist()}
- Shape: {df.shape[0]} rows × {df.shape[1]} columns
- Missing values: {df.isna().sum().sum()} total

You must answer user questions using Python code with pandas-like logic. Perform all necessary computations, filtering, and aggregation using the Python interpreter tool.

This dataset may contain missing values. If necessary, handle them appropriately (e.g., by filtering or imputing them) before analysis."""

@tool
def get_information_from_excel(file_path: str) -> str:
    """
    Extracts data and, when available, background color metadata from an Excel file (.xlsx or .xls).

    This function handles both Excel formats:
    - For .xlsx: extracts the full table and cell background colors.
    - For .xls: extracts only the tabular data (background colors not supported).

    Args:
        file_path (str): Path to the Excel file (.xlsx or .xls).

    Returns:
        str: A detailed string containing:
            - Table data extracted from the first sheet.
            - If available, a matrix of background cell colors (hex format).
            - Summary about number of rows and columns.
    """
    ext = os.path.splitext(file_path)[-1].lower()

    if ext == ".xlsx":
        # Load with openpyxl to extract both data and colors
        workbook = load_workbook(file_path, data_only=True)
        sheet = workbook.active
        color_data = []

        for row in sheet.iter_rows():
            row_colors = []
            for cell in row:
                fill = cell.fill
                color = fill.start_color.rgb if fill and fill.start_color and fill.start_color.rgb else "None"
                row_colors.append(f"#{color}" if color != "None" else "None")
            color_data.append(row_colors)

        df = pd.read_excel(file_path, engine="openpyxl", header=None)
        return f"""You are given two tables extracted from an Excel file (.xlsx):

- `data_table` contains the content of each cell.
- `color_table` contains the background color of each cell in ARGB hex format (or "None" if not set).

"data_table": {df.values.tolist()}
"color_table": {color_data}

Both tables are the same size: {df.shape[0]} rows × {df.shape[1]} columns.
Use this data to analyze the spreadsheet and answer user questions."""

    elif ext == ".xls":
        # Load with xlrd (colors not supported)
        df = pd.read_excel(file_path, engine="xlrd", header=None)
        return f"""You are given a table extracted from an Excel file (.xls):

- Background color metadata is not available for .xls files.

"data_table": {df.values.tolist()}

The table has {df.shape[0]} rows × {df.shape[1]} columns.
Use this data to analyze the spreadsheet and answer user questions."""

    else:
        return "Unsupported file format. Please provide a .xls or .xlsx Excel file."

@tool
def get_information_from_xml(file_path: str) -> str:
    """
    Reads the contents of an XML file and returns it as plain text.

    This function loads and parses an XML file from the specified file path and returns
    its content as a plain text string. The text includes the structure and data contained
    within the XML, which can help in understanding the layout and details of the document.

    Args:
        file_path (str): The local path to the XML file (e.g., "/path/to/document.xml").

    Returns:
        str: A plain text representation of the XML file's contents, including its structure and data.
    """
    with open(file_path, 'r', encoding='utf-8') as f:
        xml_string = f.read()
        xml_string = xml_string.replace(">", ">\n")
        #data_dict = xmltodict.parse(xml_string)
        #pretty_json = json.dumps(data_dict, indent=2)
    return f"""You are given an XML document converted into JSON format:

{xml_string}

Use this to answer questions about the document's structure, contents, or metadata.
If relevant, identify key entities, attributes, or relationships. """

@tool
def get_information_from_json(file_path: str) -> str:
    """
    Loads and returns JSON content as a formatted string.

    This function opens a JSON file from the specified file path, loads its contents,
    and returns the entire JSON structure as a string. The data is formatted as valid JSON,
    which is helpful for examining structured data in key-value pairs.

    Args:
        file_path (str): The local path to the .json file (e.g., "/path/to/data.json").

    Returns:
        str: A stringified version of the entire JSON data, formatted as a valid JSON string.

    Example:
        file_path = "/path/to/data.json"
        json_content = get_information_from_json(file_path)
        print(json_content)

    In this example, the function will:
    - Open the specified JSON file.
    - Load and parse its contents.
    - Return the data as a JSON-formatted string that can be used for further processing or analysis.
    """

    with open(file_path, 'r', encoding='utf-8') as f:
        data = json.load(f)
    return f"""The content of the JSON/JSON-LD file has been successfully extracted:

{json.dumps(data, indent=2)}

Instruction for the agent:

1. Scan the JSON content for any fields containing URLs (e.g., 'url', '@id', 'source', 'link').
2. For each URL found, you must explore the associated website as thoroughly as possible.
   - Do not limit the navigation to the landing page.
   - Recursively follow internal links, sections, and dynamically loaded content.
   - Extract valuable text, metadata, references, and any structured content.
3. Prioritize using the tools in the following order:

TOOL PRIORITY ORDER:
1. `navigate_browser`– Use for navigating to external websites.
2. `get_web_page` – Use for static HTML websites without JavaScript.
3. `duckduckgo_web_search_run` – Use when a URL is missing or additional context is needed.

If the JSON contains names, terms, or identifiers without direct links, consider searching Wikipedia or arXiv for background knowledge.

Important:
You must fully explore the websites found in the JSON, including internal resources and linked data. Don't stop at the home or landing page."""

@tool
def get_information_from_pdb(file_path: str) -> str:
    """
    Extracts 3D structural data from a PDB (Protein Data Bank) file.

    This function parses a PDB file and generates a human-readable summary of its
    molecular structure, including information about chains, residues, and atom positions
    in 3D space. The output is particularly useful for understanding the structural
    layout of proteins or other molecules stored in the PDB format.

    Args:
        file_path (str): The local path to the .pdb file (e.g., "/path/to/structure.pdb").

    Returns:
        str: A detailed, human-readable summary of the PDB file, listing the chains,
             residues, and atoms with their 3D coordinates.

    Example:
        file_path = "/path/to/structure.pdb"
        pdb_info = get_information_from_pdb(file_path)
        print(pdb_info)

    In this example, the function will:
    - Parse the PDB file located at the given path.
    - Extract information about chains, residues, and atoms, along with their coordinates in 3D space.
    - Return a summary of the molecular structure with the positions of atoms in x, y, and z coordinates.
    """
    parser = PDBParser(QUIET=True)
    structure = parser.get_structure("my_protein", file_path)

    info = []

    for model in structure:
        for chain in model:
            for residue in chain:
                res_info = f"Chain {chain.id} Residue {residue.resname} {residue.id[1]}:\n"
                for atom in residue:
                    coords = atom.get_coord()
                    res_info += f"  Atom {atom.name}: x={coords[0]:.2f}, y={coords[1]:.2f}, z={coords[2]:.2f}\n"
                info.append(res_info)

    return "\n".join(info)

@tool
def get_information_from_txt(file_path: str) -> str:
    """
    Extracts plain text content from a .txt file.

    This function loads the full textual content from a local .txt file and returns it
    as a single string. It preserves the original formatting and is useful for processing
    documents that contain unstructured or natural language text.

    Args:
        file_path (str): Absolute or relative path to the .txt file (e.g., "documents/file.txt").

    Returns:
        str: The complete plain text extracted from the file.

    Example:
        file_path = "notes/lecture.txt"
        text = get_information_from_txt(file_path)
        print(text)

    Use this tool when:
    - You need to read or analyze the full content of a text file.
    - The input is stored in a standard plain-text format (.txt).
    """
    loader = TextLoader(file_path)
    documents = loader.load()
    content = "\n".join([doc.page_content for doc in documents])
    return f"The TXT file content is:\n\n {content}"

@tool
def get_information_from_python(file_path: str) -> str:
    """
    Loads and returns the source code from a Python (.py) file.

    This function opens a local Python script from the specified file path, reads its entire
    source code, and returns it as a formatted string. It's useful for inspecting, analyzing,
    or executing the contents of Python files in later steps.

    Args:
        file_path (str): Absolute or relative path to the Python file (e.g., "scripts/my_script.py").

    Returns:
        str: The full Python source code from the file as a string.

    Example:
        file_path = "models/model_utils.py"
        code = get_information_from_python(file_path)
        print(code)

    Use this tool when:
    - You need to analyze or execute the contents of a Python script.
    - The input is a .py file and contains valid Python code.
    """
    with open(file_path, 'r', encoding='utf-8') as f:
        code = f.read()
    return f"""Python source code extracted from: {file_path}

{code}

Suggested next step:
If you want to analyze or run this code, use the python_code_executor tool and pass the code directly as input.
"""

@tool
def get_information_from_docx(file_path: str) -> str:
    """
    Extracts text content from a Microsoft Word (.docx) file.

    This function reads a .docx file from the specified file path and returns all
    human-readable text as a single string, preserving the logical order of the document.
    It is useful when you need to analyze or summarize documents created in Microsoft Word.

    Args:
        file_path (str): Full or relative path to the .docx file (e.g., "reports/report.docx").

    Returns:
        str: The complete extracted text from the Word document.

    Example:
        file_path = "contracts/agreement.docx"
        text = get_information_from_docx(file_path)
        print(text)

    Use this tool when:
    - The input file is a .docx Word document.
    - You want to extract the entire textual content for processing, querying, or summarization.
    """
    loader = Docx2txtLoader(file_path)
    documents = loader.load()
    return "\n".join([doc.page_content for doc in documents])

@tool
def get_information_from_pptx(file_path: str) -> str:
    """
    Extracts all text from a PowerPoint (.pptx) file slide by slide.

    Use this tool when:
    - You need to read and analyze the text content of a presentation.
    - You want to understand slide structure or extract meaningful information from slides.

    Args:
        file_path (str): Path to the .pptx file.

    Returns:
        str: Extracted text content from the presentation.
    """
    prs = Presentation(file_path)
    all_text = []

    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text.strip())
        if slide_text:
            all_text.append(f"- Slide {i + 1}\n" + "\n".join(slide_text))
    all_text = '\n\n'.join(all_text)
    return f"""The PPTX file content is:

{all_text}

You can use this information to answer the user question.
"""

@tool
def get_information_from_image(file_path: str, question: str) -> str:
    """
    Performs visual question answering (VQA) on an image file.

    This tool allows the agent to reason about the contents of an image. It takes a path to an image file
    and a natural language question, processes the image and text using a vision-language model,
    and returns a text-based answer derived from the visual information.

    Args:
        file_path (str): Path to the image file (e.g., "images/photo.png").
        question (str): Natural language question to ask about the image (e.g., "What is the person holding?").

    Returns:
        str: Answer generated by the vision-language model based on the image and the question.

    Example:
        answer = get_information_from_image("cat.jpg", "What color is the cat?")
        print(answer)

    Use this tool when:
    - You need to analyze or describe visual content in an image.
    - The user asks a question involving a photo, diagram, chart, screenshot, or other image file.
    - Visual reasoning is required to answer the question.
    """
    messages = [
        {
            "role": "user",
            "content": [
                {
                    "type": "image",
                    "image": file_path,
                },
                {"type": "text", "text": question},
            ],
        }
    ]

    text = processor_image.apply_chat_template(
        messages, tokenize=False, add_generation_prompt=True
    )
    image_inputs, video_inputs = process_vision_info(messages)
    inputs = processor_image(
        text=[text],
        images=image_inputs,
        videos=video_inputs,
        padding=True,
        return_tensors="pt",
    )
    inputs = inputs.to("cuda")

    generated_ids = model_image.generate(**inputs, max_new_tokens=1024)
    generated_ids_trimmed = [
        out_ids[len(in_ids) :] for in_ids, out_ids in zip(inputs.input_ids, generated_ids)
    ]
    output_text = processor_image.batch_decode(
        generated_ids_trimmed, skip_special_tokens=True, clean_up_tokenization_spaces=False
    )
    return f"VISUAL ANSWER: {output_text}"


@tool
def get_all_files_from_zip(file_path: str) -> Tuple[List[str], str]:
    """
    Extracts all files from a ZIP archive to a temporary directory and returns their paths.

    Args:
        file_path (str): Path to the .zip archive.

    Returns:
        Tuple[List[str], str]:
            - A list of full paths to the extracted files.
            - A summary string indicating how many files were extracted.
    """
    files = []
    temp_dir = tempfile.gettempdir()
    zip_ref = zipfile.ZipFile(file_path, 'r')
    zip_ref.extractall(temp_dir)
    zip_files = zip_ref.namelist()
    for filename in zip_ref.namelist():
        temp_file_path = os.path.join(temp_dir, filename)
        files.append(temp_file_path)

    prompt = f"The ZIP file contains {len(files)} file(s)."
    return files, prompt


duckduckgosearchrun_tool = DuckDuckGoSearchRun()
duckduckgosearchresults_tool = DuckDuckGoSearchResults(max_results=5, output_format='list')

@retry(
    wait=wait_exponential(multiplier=1, min=15, max=60),
    stop=stop_after_attempt(3),
    retry=retry_if_exception_type(DuckDuckGoSearchException)
)
@sleep_and_retry
@limits(calls=1, period=30)
def get_dgg_search(query: str, mode: str):
    if mode == "run":
        return duckduckgosearchrun_tool.run(query)
    elif mode == "results":
        return duckduckgosearchresults_tool.run(query)

def get_information_from_searches(query: str) -> str:
    """
    Perform a DuckDuckGo search and return both a textual summary and structured results.

    Args:
        query (str): The search query string.

    Returns:
        str: Combined summary and list of links with navigation instructions.
    """
    text_summary = get_dgg_search(query, mode="run")
    structured_results = get_dgg_search(query, mode="results")

    formatted_links = "\n".join([
        f"- [{res['title']}]({res['link']}):\n{res['snippet']}"
        for res in structured_results])

    text_summary = "No results found." if len(text_summary) < 5 else text_summary
    formatted_links = "No results found." if len(formatted_links) < 5 else formatted_links

    return f"""DuckDuckGo Search Results for: "{query}"

Summary:
{text_summary}

Top Links:
{formatted_links}

Instructions for the Agent:
If the summary is insufficient or lacks detail:
- Use `navigate_browser` to open one or more of the above links.
- Extract the full page content, including main text, tables, and sidebars.
- Optionally, explore internal links within those pages if relevant to the query.
"""

search_tool = Tool(
    name="duckduckgo_search",
    func=get_information_from_searches,
    description=(
        "Use this tool to perform a live DuckDuckGo web search and retrieve both a summary and a list of relevant links "
        "with snippets. Ideal for finding general information from public websites that are well indexed by DuckDuckGo.\n\n"
        "Limitations:\n"
        "- DuckDuckGo may not return results for sites like Google.\n"
        "Fallback Instructions:\n"
        "If this tool returns no useful information or fails to retrieve results:\n"
        "- Use `navigate_browser` to open relevant links or perform a direct search in a browser context.\n"
        "- Extract detailed information from the visited web pages, including main text, sidebars, and tables.\n"
        "- Explore internal links within the pages if needed to answer the question."
    )
)

# Python interpreter
def verbose_python_executor(code: str) -> str:
    """Executes Python code and returns both the code and the output."""
    tool = PythonREPLTool()
    result = tool.run(code)
    return f"Code executed:\n```python\n{code}\n```\n\nOutput:\n{result}"

python_tool = Tool(
    name="python_code_executor",
    func=verbose_python_executor,
    description=(
        "Use this tool to execute raw Python code.\n\n"
        "Input: A Python code snippet as a single string.\n\n"
        "Important:\n"
        "- This is NOT a remote API call. Do NOT include calls to `default_api`, `__arg1`, or any tool inside the code.\n"
        "- The code should be ready to run directly in Python, using standard libraries."
    )
)

def download_file_temp(url: str) -> str:
    """
    Downloads a file from the given URL and saves it into a temporary directory.

    Args:
        url (str): URL of the file to download.

    Returns:
        str: Path to the downloaded file or an error message.
    """
    try:
        response = requests.get(url, stream=True)
        response.raise_for_status()

        # Create a temporary directory
        temp_dir = tempfile.mkdtemp()

        # Extract filename from URL or fallback to a generic name
        filename = url.split("/")[-1] or "downloaded_file"

        # Full path for saving file
        file_path = os.path.join(temp_dir, filename)

        with open(file_path, "wb") as f:
            for chunk in response.iter_content(chunk_size=8192):
                f.write(chunk)

        return f"File downloaded successfully and saved at {file_path}"
    except Exception as e:
        return f"Failed to download file: {e}"

download_tool = Tool(
    name="file_downloader_temp",
    func=download_file_temp,
    description="Downloads a file from a URL and saves it in a temporary folder."
)

# Create a Playwright browser instance
def initialize_web_tools():
    async_browser = create_async_playwright_browser()
    toolkit = PlayWrightBrowserToolkit.from_browser(async_browser=async_browser)
    return toolkit.get_tools()
